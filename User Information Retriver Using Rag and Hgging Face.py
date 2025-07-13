# rag_autonomous_retriever.py
# Complete RAG System: Ingest Any File Format, Summarize, and Answer with Citations

import os
import json
import pandas as pd
import numpy as np
import pdfplumber
import docx
import openpyxl
import pptx
import faiss
import glob
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from flask import Flask, request, jsonify, render_template

app = Flask(__name__)
UPLOAD_FOLDER = 'uploaded_files'
EMBEDDING_DIM = 384
model = SentenceTransformer('all-MiniLM-L6-v2')

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Global Indexes
text_chunks = []
chunk_sources = []
index = faiss.IndexFlatL2(EMBEDDING_DIM)


def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    elif ext == '.pdf':
        with pdfplumber.open(filepath) as pdf:
            return "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    elif ext == '.docx':
        doc = docx.Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == '.csv':
        df = pd.read_csv(filepath)
        return df.to_string()
    elif ext == '.xlsx':
        df = pd.read_excel(filepath)
        return df.to_string()
    elif ext == '.pptx':
        prs = pptx.Presentation(filepath)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return ""


def chunk_text(text, max_words=100):
    words = text.split()
    return [" ".join(words[i:i + max_words]) for i in range(0, len(words), max_words)]


def index_uploaded_files():
    global text_chunks, chunk_sources, index
    text_chunks.clear()
    chunk_sources.clear()
    index.reset()
    files = glob.glob(os.path.join(UPLOAD_FOLDER, '*'))
    for file in files:
        raw_text = extract_text_from_file(file)
        chunks = chunk_text(raw_text)
        embeddings = model.encode(chunks)
        text_chunks.extend(chunks)
        chunk_sources.extend([f"{os.path.basename(file)} (chunk {i+1})" for i in range(len(chunks))])
        index.add(np.array(embeddings))


def answer_question(query, top_k=3):
    q_embedding = model.encode([query])
    D, I = index.search(np.array(q_embedding), top_k)
    relevant_chunks = [text_chunks[i] for i in I[0]]
    citations = [chunk_sources[i] for i in I[0]]
    answer = f"Answer based on retrieved chunks:\n\n"
    for i in range(len(relevant_chunks)):
        answer += f"[{i+1}] {citations[i]}: {relevant_chunks[i]}\n\n"
    return answer

@app.route('/')
def upload_page():
    return render_template('upload.html')

@app.route('/upload', methods=['POST'])
def upload_files():
    if 'files[]' not in request.files:
        return "No files uploaded", 400
    files = request.files.getlist('files[]')
    for file in files:
        filepath = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(filepath)
    index_uploaded_files()
    return "Files uploaded and indexed successfully."

@app.route('/ask', methods=['POST'])
def ask():
    data = request.get_json()
    query = data.get('question')
    if not query:
        return jsonify({"error": "No question provided"}), 400
    result = answer_question(query)
    return jsonify({"answer": result})

if __name__ == '__main__':
    index_uploaded_files()
    app.run(debug=True)
